# Generic AI Document-Based Content Generator for Any Business

import streamlit as st
import google.generativeai as genai
import os
from docx import Document
from pptx import Presentation
import PyPDF2
import pandas as pd
import tempfile

api_key = "AIzaSyC9WS0oHMIaFCbgqxI-gYzNNwG9rjxRbIk"
genai.configure(api_key=api_key)

REQUIRED_DOCS = ['brochure', 'email_ppt', 'catalog', 'email', 'price_sheet']
TOOL_MAP = {
    'sales_pitch': 'Sales Pitch Generator',
    'email': 'Email Template Generator',
    'linkedin': 'LinkedIn Message Generator',
    'whatsapp': 'WhatsApp Message Generator',
    'proposal': 'Proposal Generator',
    'product_explanation': 'Product Explanation',
    'technical_summary': 'Technical Summary',
    'resume_screener': 'Resume Screener & JD Creator'
}

# --- Prompt Map ---
def get_prompt(mode):
    prompt_map = {
        "sales_pitch": """
You are a senior B2B sales strategist with expertise in SaaS, automation, and solution-based selling.

Carefully review and analyze all the uploaded business documents (brochure, PPT, catalogue, email, and pricing sheet). Based strictly on the content of these documents:

Craft a clear, compelling, and concise B2B spoken-style sales pitch that can be delivered in 30–40 seconds.

Your pitch must:
1. Start with a confident self-introduction and company name
2. Clearly state what problem the product/service solves for the target customer
3. Briefly highlight 2–3 high-impact features or differentiators
4. Mention a real or representative result, outcome, or credibility point (e.g., customer success, ROI, trust)
5. End with a friendly and professional call-to-action that invites a follow-up conversation

Formatting guidelines:
- Write the output as a natural spoken pitch – avoid headings or bullet points
- Keep the tone professional, persuasive, and human – not robotic or generic
- Do not use markdown, symbols, emojis, or asterisks
- Keep the total length within 4–6 sentences (~120–150 words)
- Make the content specific to the uploaded documents

Do not add any disclaimers or generic filler — focus on delivering a personalized pitch that a sales expert would use in a real conversation.
""",
        "linkedin": """
You are a senior B2B sales and networking strategist experienced in SaaS and enterprise outreach on LinkedIn.

Carefully analyze all the uploaded business documents (brochure, PPT, catalogue, email, pricing sheet). Based strictly on the content of these documents:

Write a short, professional LinkedIn outreach message to a potential customer or decision-maker.

The message must:
1. Open with a natural greeting and a quick, non-generic reason for reaching out
2. Briefly describe what the company does and the specific problem it solves
3. Mention 1–2 impressive features, results, or differentiators supported by the documents
4. Close with a soft, professional call-to-action (like offering to share insights, case studies, or hop on a short call)

Formatting rules:
- Total message length should be 50–70 words (2–4 short sentences)
- Avoid emojis, asterisks, hashtags, markdown, or robotic phrases
- Keep it conversational and personalized, like a real LinkedIn message — not a mass email
- Do not include subject lines, sign-offs, or email-style formatting

This message should feel like it’s coming from a helpful professional — not a bot or aggressive salesperson.
""",
        "email": """
You are a senior B2B marketing and sales strategist with deep experience in SaaS, automation, and enterprise outreach.

Carefully review and analyze all the uploaded business documents (brochure, PPT, catalogue, email threads, and pricing sheet). Based strictly on the content of these documents:

Write a concise, high-converting B2B email that introduces the product or service to a potential customer.

Your email must follow this structure:
1. Personalized and attention-grabbing opening line (related to the recipient’s likely challenge or industry trend)
2. Clear and compelling value proposition (what the product/service does and how it helps)
3. 2–3 specific features, differentiators, or outcomes supported by the documents
4. A credibility line (customer results, market trust, certifications, etc.)
5. A friendly, low-friction call-to-action (like suggesting a short call or offering more details)

Formatting rules:
- Write a complete email with Subject line, Greeting, Body, CTA, and Signature
- Keep the tone professional, helpful, and human — not robotic or overly generic
- Do not use asterisks, emojis, markdown symbols, or boilerplate phrases
- Keep the total email length between 100–140 words
- Avoid fluff — focus on real, document-backed value

The output should read like a real sales email from a smart, helpful sales executive — not like an AI or a mass marketing bot.
""",
        "whatsapp": """
You are a B2B communication expert specializing in WhatsApp marketing.

Carefully review all uploaded business documents (brochure, PPT, emails, and pricing). Based strictly on the content:

Generate a short, engaging WhatsApp message that introduces the product or service to a potential lead.

Guidelines:
- Keep the message under 50 words
- Tone: Friendly, crisp, and professional
- Focus on one clear benefit or outcome (e.g., cost-saving, automation, performance)
- End with a soft call-to-action (e.g., “Interested in a quick chat?” or “Want more details?”)

Avoid emojis, excessive punctuation, or robotic phrases. This should feel like a real message from a sales rep to a qualified lead.
""",
        "proposal": """
You are a senior business proposal writer with experience in B2B SaaS, automation, and enterprise solutions.

Review all uploaded documents (brochure, PPT, catalogue, pricing sheet, emails). Based strictly on their content:

Create a well-structured business proposal that includes:
1. Title
2. Executive Summary (summarize client pain points + your solution)
3. Solution Overview (features, benefits, and differentiators)
4. Impact/Results (case studies, outcomes, or stats)
5. Pricing or Package Overview (based on available data)
6. Next Steps (suggest a call or demo)

Use formal formatting with section headers. Do not use markdown, emojis, or casual language. Keep it business-focused and aligned to the customer's potential goals.
""",
        "product_explanation": """
You are a senior product marketing specialist.

Read all uploaded business documents (brochure, catalogue, PPT, and pricing). Based strictly on their content:

Write a clear, engaging explanation of the product or service for a professional audience unfamiliar with it.

Guidelines:
- Keep it concise: 150–200 words
- Focus on what the product does, who it's for, and why it matters
- Highlight key features, differentiators, and outcomes
- Use professional language with short, informative sentences

Avoid technical jargon unless it’s explained. Do not use bullet points, emojis, or marketing fluff.
""",
        "technical_summary": """
You are a technical documentation expert with experience in SaaS and automation platforms.

Analyze all uploaded documents (technical brochure, architecture slides, catalogue, feature list).

Generate a technical summary of the product that includes:
1. Core functionality and how it works
2. Key architecture or technical components (if available)
3. Integrations, automation, or platform compatibility
4. Performance, scalability, or security benefits

Keep the summary under 300 words. Use clear, structured language for technical readers (e.g., engineers, IT managers). Avoid promotional tone, markdown, emojis, or fluff.
""",
        "resume_screener": """
You are an HR tech assistant designed to analyze candidate resumes and job descriptions for alignment.

Using the uploaded resume and job description:
- Highlight key skill and experience matches
- Identify any major gaps
- Score alignment on a scale from 1 to 10
- Provide a 3-line summary of how well the candidate fits the role

Keep the tone professional, neutral, and helpful. Do not use emojis, marketing tone, or unnecessary commentary.
"""
    }
    return prompt_map.get(mode, "You are an AI assistant. Read all uploaded business documents and summarize the business context professionally.")

# --- Helper Functions ---
def extract_text_from_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.txt':
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    elif ext == '.docx':
        doc = Document(filepath)
        return '\n'.join([para.text for para in doc.paragraphs])
    elif ext == '.pptx':
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text.strip())
        return '\n'.join(text)
    elif ext == '.pdf':
        text = []
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
        return '\n'.join(text)
    elif ext == '.csv':
        df = pd.read_csv(filepath)
        return df.to_string(index=False)
    elif ext == '.xlsx':
        df = pd.read_excel(filepath)
        return df.to_string(index=False)
    else:
        return ''

def generate_content(mode, data):
    intro = get_prompt(mode)
    prompt = f"""{intro}

--- START OF MATERIAL ---
{data}
--- END OF MATERIAL ---

Ensure the output is concise, professional, and formatted cleanly. Do not use asterisks, hashtags, or emojis.
"""
    model = genai.GenerativeModel('gemini-2.5-pro')
    response = model.generate_content(prompt)
    return response.text.strip()

# --- Streamlit App ---
st.set_page_config(page_title="AI Document Content Generator", layout="wide")
st.title("AI Document-Based Content Generator")

if 'uploaded_files' not in st.session_state:
    st.session_state.uploaded_files = {doc: None for doc in REQUIRED_DOCS}
if 'result' not in st.session_state:
    st.session_state.result = None
if 'selected_tool' not in st.session_state:
    st.session_state.selected_tool = list(TOOL_MAP.keys())[0]

with st.sidebar:
    st.header("1. Upload Business Documents (any subset)")
    for doc in REQUIRED_DOCS:
        uploaded = st.file_uploader(f"Upload {doc.replace('_', ' ').title()}", type=["pdf", "docx", "pptx", "txt", "csv", "xlsx"], key=doc)
        if uploaded:
            st.session_state.uploaded_files[doc] = uploaded
    if st.button("Clear All Uploads"):
        st.session_state.uploaded_files = {doc: None for doc in REQUIRED_DOCS}
        st.session_state.result = None
        st.success("All uploads cleared.")
    st.header("2. Select Tool")
    st.session_state.selected_tool = st.selectbox("Choose a content generator tool:", list(TOOL_MAP.keys()), format_func=lambda x: TOOL_MAP[x])

st.subheader("Upload Status:")
any_uploaded = False
for doc in REQUIRED_DOCS:
    if st.session_state.uploaded_files[doc]:
        st.write(f"✅ {doc.replace('_', ' ').title()} uploaded.")
        any_uploaded = True
    else:
        st.write(f"❌ {doc.replace('_', ' ').title()} not uploaded.")

if any_uploaded:
    if st.button(f"Generate {TOOL_MAP[st.session_state.selected_tool]}"):
        # Save files to temp dir and extract text
        docs_content = []
        with tempfile.TemporaryDirectory() as tmpdir:
            for doc in REQUIRED_DOCS:
                uploaded = st.session_state.uploaded_files[doc]
                if uploaded:
                    ext = os.path.splitext(uploaded.name)[1].lower()
                    temp_path = os.path.join(tmpdir, doc + ext)
                    with open(temp_path, 'wb') as f:
                        f.write(uploaded.read())
                    docs_content.append(extract_text_from_file(temp_path))
            data = '\n'.join(docs_content)
            with st.spinner("Generating content..."):
                try:
                    result = generate_content(st.session_state.selected_tool, data)
                    st.session_state.result = result
                except Exception as e:
                    st.error(f"Error generating content: {e}")
else:
    st.info("Please upload at least one business document to enable content generation.")

if st.session_state.result:
    st.subheader(f"{TOOL_MAP[st.session_state.selected_tool]} Result:")
    st.write(st.session_state.result)
